#!/usr/bin/env python3
"""
Generate OutYah investor deck — cloned from Find-your-next-outing.pptx layout,
with copy tuned for investors and live catalog screenshots.
"""

from __future__ import annotations

import io
import json
import shutil
import tempfile
import urllib.error
import urllib.request
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path(__file__).resolve().parents[1]
ASSETS = Path(__file__).resolve().parent / "deck_assets"
TEMPLATE = ASSETS / "Find-your-next-outing.pptx"
REF_DOWNLOADS = Path.home() / "Downloads" / "Find-your-next-outing.pptx"
OUT = Path.home() / "Downloads" / "OutYah-Investor-Deck.pptx"
LIVE = "vtdi-platform.vercel.app"


def ensure_template() -> Path:
    if REF_DOWNLOADS.exists():
        ASSETS.mkdir(parents=True, exist_ok=True)
        if not TEMPLATE.exists() or REF_DOWNLOADS.stat().st_mtime > TEMPLATE.stat().st_mtime:
            shutil.copy2(REF_DOWNLOADS, TEMPLATE)
        return TEMPLATE
    if TEMPLATE.exists():
        return TEMPLATE
    raise FileNotFoundError(
        "Missing template. Place Find-your-next-outing.pptx in Downloads or scripts/deck_assets/"
    )


def http_get(url: str, *, headers: dict | None = None, timeout: int = 30) -> bytes | None:
    req = urllib.request.Request(url, headers=headers or {"User-Agent": "OutYahDeck/1.0"})
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return resp.read()
    except (urllib.error.URLError, TimeoutError, OSError):
        return None


def fetch_catalog_images() -> list[str]:
    env = {}
    env_path = ROOT / ".env"
    if env_path.exists():
        for line in env_path.read_text().splitlines():
            if "=" in line and not line.strip().startswith("#"):
                k, _, v = line.partition("=")
                v = v.strip().strip('"').strip("'")
                env[k.strip()] = v
    url = env.get("VITE_SUPABASE_URL")
    key = env.get("VITE_SUPABASE_ANON_KEY")
    if not url or not key:
        return []
    headers = {"apikey": key, "Authorization": f"Bearer {key}"}
    raw = http_get(
        f"{url}/rest/v1/places?select=image&image=not.is.null&limit=6",
        headers=headers,
        timeout=20,
    )
    if not raw:
        return []
    rows = json.loads(raw.decode("utf-8"))
    return [r["image"] for r in rows if r.get("image")]


def download(url: str, cache: Path) -> Path | None:
    if not url:
        return None
    name = url.split("/")[-1].split("?")[0] or "img.jpg"
    dest = cache / name
    if dest.exists():
        return dest
    data = http_get(url)
    if not data:
        return None
    dest.write_bytes(data)
    return dest


def set_text(shape, text: str) -> None:
    tf = shape.text_frame
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if p.runs:
            p.runs[0].text = line
            for run in p.runs[1:]:
                run.text = ""
        else:
            p.text = line
        for _ in range(len(tf.paragraphs) - len(lines)):
            tf._element.remove(tf.paragraphs[-1]._p)


def text_shapes(slide):
    return [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]


def replace_picture(slide, shape, path: Path):
    left, top, width, height = shape.left, shape.top, shape.width, shape.height
    with Image.open(path) as im:
        im = im.convert("RGB")
        buf = io.BytesIO()
        im.save(buf, format="JPEG", quality=92)
        buf.seek(0)
    sp = shape._element
    sp.getparent().remove(sp)
    slide.shapes.add_picture(buf, left, top, width=width, height=height)


def apply_copy(prs: Presentation) -> None:
    # Slide 1 — title (match reference tone)
    s1 = text_shapes(prs.slides[0])
    set_text(s1[0], "OutYah")
    set_text(s1[1], "Find your next outing.")
    set_text(
        s1[2],
        "A Jamaica-wide discovery platform — 14 parishes, curated venues, real plans. Built by Terobytez.",
    )
    set_text(s1[3], "🚀 LIVE NOW")
    set_text(s1[4], LIVE)

    # Slide 2 — problem (blend Find-your-next-outing + Untitled)
    s2 = text_shapes(prs.slides[1])
    set_text(s2[0], "The Problem We're Solving")
    set_text(s2[1], "Jamaica's outing scene is fragmented")
    set_text(
        s2[2],
        "Events, venues, and experiences live across Instagram, WhatsApp groups, and word-of-mouth. "
        "There's no single trusted source for where you should go tonight.",
    )
    set_text(
        s2[3],
        "No centralized discovery tool for Jamaican outings\n"
        "Hard to verify quality before you go\n"
        "No easy way to plan, share, and navigate a full outing",
    )

    # Slide 3 — unchanged from reference (Discover / Trust / Plan / Go)
    s3 = text_shapes(prs.slides[2])
    set_text(s3[0], "Introducing OutYah")
    set_text(s3[1], "One platform. Every outing. All of Jamaica.")

    # Slide 4 — data model
    s4 = text_shapes(prs.slides[3])
    set_text(s4[0], "How It's Built — The Data Model")
    set_text(s4[1], "Eight core objects power every experience on OutYah.")
    set_text(
        s4[2],
        "At the heart of OutYah is the Post — it connects a User to a Venue, tagged with a "
        "Category and Parish for discovery.\n"
        "Users build OutingPlans made of PlanItems, and save favorites across the platform.",
    )

    # Slide 5 — user journey
    s5 = text_shapes(prs.slides[4])
    set_text(s5[0], "The User Journey")
    set_text(s5[1], "Five steps from discovery to showing up.")

    # Slide 6 — stack
    s6 = text_shapes(prs.slides[5])
    set_text(s6[0], "Built On")
    set_text(s6[1], "A modern, production-grade stack — built to scale.")

    # Slide 7 — live features (investor refresh)
    s7 = text_shapes(prs.slides[6])
    set_text(s7[0], "Live Today — Shipped, Not a Concept")
    set_text(
        s7[1],
        f"OutYah is a real product with real users. Demo it now at {LIVE}.",
    )
    # keep first four feature blocks; refresh last two for shipped wow features
    set_text(s7[2], "✅ Auth & Profiles")
    set_text(s7[3], "Full user authentication with personalized profiles")
    set_text(s7[4], "✅ Parish & Category Filtering")
    set_text(s7[5], "Map-based explore across all 14 parishes")
    set_text(s7[6], "✅ Curated Listings")
    set_text(s7[7], "Booking info, image-first posts, and admin approval queue")
    set_text(s7[8], "✅ Outing Planner")
    set_text(s7[9], "Build plans, reorder stops, share via link, open in Google Maps")
    set_text(s7[10], "✅ Jamaica Pulse")
    set_text(s7[11], "Live strip — on now, open venues, weather, and what's on tonight")
    set_text(s7[12], "✅ JMD Costs & Event Chat")
    set_text(s7[13], "Outing cost estimates in JMD, RSVP counts, and gated realtime chat")

    # Slide 8 — product in action
    s8 = text_shapes(prs.slides[7])
    set_text(s8[0], "Product in Action")
    set_text(
        s8[1],
        "From the branded landing page to the admin portal — every surface is designed, shipped, and live.",
    )

    # Slide 9 — why OutYah wins
    s9 = text_shapes(prs.slides[8])
    set_text(s9[0], "Why OutYah Wins")
    set_text(s9[1], "Built for Jamaica, by Jamaicans")
    set_text(
        s9[2],
        "OutYah understands local context — parish culture, event types, and how people actually discover outings here.",
    )
    set_text(s9[3], "Trust-first — curated listings with real reviews, not noisy map pins")
    set_text(s9[4], "Shareable by design — plans spread organically through WhatsApp")
    set_text(s9[5], "Already live — investors can demo the product tonight, not next quarter")

    # Slide 10 — vision
    s10 = text_shapes(prs.slides[9])
    set_text(s10[0], "The Vision")
    set_text(s10[1], "Make it effortless to trust where you go in Jamaica.")
    set_text(
        s10[2],
        "OutYah is just getting started. The foundation is built. The platform is live. The community is growing.",
    )
    set_text(
        s10[3],
        f"🌐 See it live: {LIVE} — built by Terobytez for VTDI/UTech.",
    )


def swap_product_images(prs: Presentation, cache: Path) -> None:
    """Replace slide 8 catalog grid with live Supabase photos when available."""
    urls = fetch_catalog_images()
    if not urls:
        return

    paths = []
    for url in urls[:6]:
        p = download(url, cache)
        if p:
            paths.append(p)
    if not paths:
        return

    slide = prs.slides[7]
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    # reference slide 8 has 6 large screenshots in top row + partial second row
    for pic, path in zip(pics[: len(paths)], paths):
        replace_picture(slide, pic, path)


def main():
    template = ensure_template()
    cache = Path(tempfile.mkdtemp(prefix="outyah-deck-"))

    shutil.copy2(template, OUT)
    prs = Presentation(str(OUT))
    apply_copy(prs)
    swap_product_images(prs, cache)
    prs.save(str(OUT))

    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Template: {template}")


if __name__ == "__main__":
    main()
